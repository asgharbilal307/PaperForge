import sys
import re
import io
import uuid
import hashlib
import logging
import argparse
import json
from pathlib import Path
from dataclasses import dataclass
from typing import Optional

import fitz
from docx import Document as DocxDocument
from pptx import Presentation
import nbformat
from github import Github, GithubException
from qdrant_client.models import PointStruct, Filter, FieldCondition, MatchValue

sys.path.insert(0, str(Path(__file__).parent.parent))

from app.core.config import get_settings
from app.core.vectorstore import get_qdrant_client, embed_texts, ensure_collection_exists

logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")
logger = logging.getLogger(__name__)

CHUNK_SIZE = 600
CHUNK_OVERLAP = 100
BATCH_SIZE = 64
STATE_FILE = Path(__file__).parent.parent / ".ingest_state.json"

DOC_TYPE_PATTERNS = [
    (r"(?<![a-z])quiz(?:zes?)?\d*(?![a-z])", "quiz"),
    (r"(?<![a-z])mid(?:term)?s?[\s\-_]?\d*(?![a-z])", "mid"),
    (r"(?<![a-z])final(?:s|[\s_\-]?exam)?(?![a-z])", "final"),
    (r"(?<![a-z])(?:notes?|lectures?|slides?|handouts?)(?![a-z])",         "notes"),
    (r"(?<![a-z])(?:assignments?|labs?|projects?|homework|hw)\d*(?![a-z])", "assignment"),
]

NON_COURSE_PATTERNS = [
    r"^[-=*]+",
    r"^\d{4}$",
    r"^(fall|spring|summer|winter)\s*\d*$",
    r"^semester\s*\d+$",
]

DOC_TYPE_WORDS = {
    "quiz","quizzes","mid","mids","midterm","midterms","final","finals",
    "notes","note","lecture","lectures","slides","slide","book","books",
    "handout","handouts","assignment","assignments","lab","labs",
    "project","projects","homework","hw","other",
}

@dataclass
class ParsedDocument:
    text: str
    course: str
    doc_type: str
    filename: str
    year: Optional[str]
    source_url: str
    raw_path: str
    full_path: str
    sha: str


# ================= STATE =================

def load_state() -> dict:
    if STATE_FILE.exists():
        try:
            return json.loads(STATE_FILE.read_text())
        except Exception:
            pass
    return {}

def save_state(state: dict):
    STATE_FILE.write_text(json.dumps(state, indent=2))


# ================= METADATA =================

def extract_year(path: str) -> Optional[str]:
    m = re.search(r"(20\d{2})", path)
    return m.group(1) if m else None

def detect_doc_type(path: str) -> str:
    lower = path.lower()
    for pattern, doc_type in DOC_TYPE_PATTERNS:
        if re.search(pattern, lower):
            return doc_type
    return "other"

def is_non_course(part: str) -> bool:
    p = part.strip().lower()
    if not p or len(p) < 2:
        return True
    if any(re.search(pat, p) for pat in NON_COURSE_PATTERNS):
        return True
    if p in DOC_TYPE_WORDS:
        return True
    if re.match(r"^\d+$", p):
        return True
    return False

def extract_course(parts: list[str], full_path: str) -> str:
    for part in parts[:-1]:
        if not is_non_course(part):
            return part.strip().lower()

    stem = Path(full_path).stem
    cleaned = re.sub(
        r"[_\-\s]+(quiz|mid|midterm|final|notes|assignment|lab|exam)\d*$",
        "", stem, flags=re.IGNORECASE
    ).strip()
    cleaned = re.sub(r"[_\-\s]+\d{4}$", "", cleaned).strip()
    return cleaned if cleaned else stem if stem else "General".lower()


# ================= PARSERS =================

def parse_pdf(b):
    try:
        doc = fitz.open(stream=b, filetype="pdf")
        return "\n".join(p.get_text() for p in doc)
    except:
        return ""

def parse_docx(b):
    try:
        doc = DocxDocument(io.BytesIO(b))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except:
        return ""

def parse_pptx(b):
    try:
        prs = Presentation(io.BytesIO(b))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)
    except:
        return ""

def parse_ipynb(b):
    try:
        nb = nbformat.reads(b.decode("utf-8", errors="ignore"), as_version=4)
        return "\n".join(cell.source for cell in nb.cells if cell.cell_type in ("markdown","code"))
    except:
        return ""


# ================= CHUNKING =================

def chunk_text(text: str):
    chunks, start = [], 0
    while start < len(text):
        chunks.append(text[start:start+CHUNK_SIZE].strip())
        start += CHUNK_SIZE - CHUNK_OVERLAP
    return [c for c in chunks if len(c) > 50]


# ================= QDRANT =================

def delete_chunks_for_path(raw_path: str):
    settings = get_settings()
    client = get_qdrant_client()
    client.delete(
        collection_name=settings.qdrant_collection,
        points_selector=Filter(
            must=[FieldCondition(key="raw_path", match=MatchValue(value=raw_path))]
        )
    )

def upsert_document(doc: ParsedDocument):
    settings = get_settings()
    client = get_qdrant_client()

    if not doc.text.strip():
        return

    chunks = chunk_text(doc.text)
    if not chunks:
        return

    all_points = []

    for i, chunk in enumerate(chunks):
        point_id = str(uuid.UUID(hashlib.md5(f"{doc.raw_path}::{i}".encode()).hexdigest()))
        all_points.append({
            "id": point_id,
            "text": chunk,
            "payload": {
                "course": doc.course,
                "doc_type": doc.doc_type,
                "filename": doc.filename,
                "year": doc.year,
                "source_url": doc.source_url,
                "raw_path": doc.raw_path,
                "full_path": doc.full_path,
                "sha": doc.sha,
                "chunk_index": i,
                "total_chunks": len(chunks),
            },
        })

    for i in range(0, len(all_points), BATCH_SIZE):
        batch = all_points[i:i+BATCH_SIZE]
        texts = [p["text"] for p in batch]
        vectors = embed_texts(texts)

        points = [
            PointStruct(
                id=batch[j]["id"],
                vector=vectors[j],
                payload={**batch[j]["payload"], "content": texts[j]},
            )
            for j in range(len(batch))
        ]

        client.upsert(collection_name=settings.qdrant_collection, points=points)

    logger.info(f"+ [{doc.doc_type}] {doc.course} | {doc.filename} => {len(chunks)} chunks")


# ================= MAIN =================

def fetch_and_ingest(force=False):
    settings = get_settings()
    gh = Github(settings.github_token)
    repo = gh.get_repo(settings.github_repo)

    ensure_collection_exists()

    state = {} if force else load_state()
    new_state = {}
    stats = {"new":0,"updated":0,"skipped":0,"failed":0}

    def walk(path=""):
        try:
            contents = repo.get_contents(path)
        except GithubException:
            return

        for item in contents:
            if item.type == "dir":
                walk(item.path)
                continue

            ext = Path(item.name).suffix.lower()
            if ext not in {".pdf",".docx",".pptx",".ipynb",".cpp",".py",".md",".txt"}:
                continue

            sha = item.sha
            new_state[item.path] = sha

            if not force and state.get(item.path) == sha:
                stats["skipped"] += 1
                continue

            try:
                raw = item.decoded_content
            except:
                stats["failed"] += 1
                continue

            if ext == ".pdf":
                text = parse_pdf(raw)
            elif ext == ".docx":
                text = parse_docx(raw)
            elif ext == ".pptx":
                text = parse_pptx(raw)
            elif ext == ".ipynb":
                text = parse_ipynb(raw)
            else:
                text = raw.decode("utf-8", errors="ignore")

            if not text.strip():
                stats["failed"] += 1
                continue

            parts = item.path.split("/")

            doc = ParsedDocument(
                text=text,
                course=extract_course(parts, item.path),
                doc_type=detect_doc_type(item.path),
                filename=item.name,
                year=extract_year(item.path),
                source_url=item.html_url,
                raw_path=item.path,
                full_path=item.path,
                sha=sha,
            )

            if item.path in state:
                delete_chunks_for_path(item.path)
                stats["updated"] += 1
            else:
                stats["new"] += 1

            upsert_document(doc)

    walk()

    for path in set(state) - set(new_state):
        delete_chunks_for_path(path)

    save_state(new_state)

    logger.info(f"\nDone: {stats}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--force", action="store_true")
    args = parser.parse_args()
    fetch_and_ingest(force=args.force)